# -*- coding: utf-8 -*-
"""内置技能初始化 — 系统预置的 5 个基础技能（seed data）

在应用首次启动时自动插入数据库，使用 INSERT IGNORE / ON DUPLICATE 语义避免重复。
"""

from __future__ import annotations

import logging
from datetime import datetime

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from ..models.skill import Skill

logger = logging.getLogger("ai_hubs.skills")

# ═══════════════════════════════════════════════════════════
# 5 个内置技能定义
# ═══════════════════════════════════════════════════════════

_BUILTIN_SKILLS: list[dict] = [
    {
        "name": "docx",
        "description": (
            "Word 文档处理技能。支持读取、创建、编辑 .docx 文件，包括格式化文本、"
            "表格、图片、页眉页脚、样式修改、批注与修订追踪。"
        ),
        "category": "document",
        "version": "1.0.0",
        "config": {
            "entry": "skill.py",
            "code": (
                "# docx skill — Word 文档处理\n"
                "# 依赖: python-docx\n"
                "# 安装: pip install python-docx\n\n"
                "from docx import Document\n"
                "from docx.shared import Inches, Pt, RGBColor\n"
                "from docx.enum.text import WD_ALIGN_PARAGRAPH\n\n"
                "def create_doc(path: str):\n"
                "    doc = Document()\n"
                "    doc.save(path)\n"
                "    return doc\n\n"
                "def read_doc(path: str) -> str:\n"
                "    doc = Document(path)\n"
                "    return '\\n'.join(p.text for p in doc.paragraphs)\n\n"
                "def add_heading(doc, text: str, level: int = 1):\n"
                "    doc.add_heading(text, level=level)\n\n"
                "def add_paragraph(doc, text: str, bold: bool = False):\n"
                "    p = doc.add_paragraph(text)\n"
                "    if bold:\n"
                "        for run in p.runs:\n"
                "            run.bold = True\n"
                "    return p\n\n"
                "def add_table(doc, rows: int, cols: int, data: list = None):\n"
                "    table = doc.add_table(rows=rows, cols=cols, style='Table Grid')\n"
                "    if data:\n"
                "        for i, row_data in enumerate(data):\n"
                "            for j, cell_text in enumerate(row_data):\n"
                "                table.cell(i, j).text = str(cell_text)\n"
                "    return table\n\n"
                "def add_image(doc, image_path: str, width: float = 4.0):\n"
                "    doc.add_picture(image_path, width=Inches(width))\n"
            ),
        },
    },
    {
        "name": "xlsx",
        "description": (
            "Excel 表格处理技能。支持读取、创建、编辑 .xlsx/.xlsm 文件，包括公式计算、"
            "数据透视表、图表、条件格式、单元格样式、合并拆分、筛选排序。"
        ),
        "category": "document",
        "version": "1.0.0",
        "config": {
            "entry": "skill.py",
            "code": (
                "# xlsx skill — Excel 表格处理\n"
                "# 依赖: openpyxl\n"
                "# 安装: pip install openpyxl\n\n"
                "from openpyxl import Workbook, load_workbook\n"
                "from openpyxl.styles import Font, PatternFill, Alignment, Border, Side\n"
                "from openpyxl.chart import BarChart, LineChart, PieChart, Reference\n"
                "from openpyxl.utils import get_column_letter\n\n"
                "def create_workbook():\n"
                "    return Workbook()\n\n"
                "def read_workbook(path: str):\n"
                "    return load_workbook(path, data_only=False)\n\n"
                "def get_sheet_names(wb) -> list:\n"
                "    return wb.sheetnames\n\n"
                "def read_sheet(wb, sheet_name: str = None):\n"
                "    ws = wb[sheet_name] if sheet_name else wb.active\n"
                "    data = []\n"
                "    for row in ws.iter_rows(values_only=True):\n"
                "        data.append(list(row))\n"
                "    return data\n\n"
                "def write_cell(ws, row: int, col: int, value, bold: bool = False):\n"
                "    cell = ws.cell(row=row, column=col, value=value)\n"
                "    if bold:\n"
                "        cell.font = Font(bold=True)\n"
                "    return cell\n\n"
                "def add_chart(ws, chart_type: str, data_start: str, data_end: str, position: str):\n"
                "    data = Reference(ws, range_string=f'{ws.title}!{data_start}:{data_end}')\n"
                "    if chart_type == 'bar':\n"
                "        chart = BarChart()\n"
                "    elif chart_type == 'line':\n"
                "        chart = LineChart()\n"
                "    elif chart_type == 'pie':\n"
                "        chart = PieChart()\n"
                "    else:\n"
                "        chart = BarChart()\n"
                "    chart.add_data(data)\n"
                "    ws.add_chart(chart, position)\n"
                "    return chart\n\n"
                "def auto_fit_columns(ws):\n"
                "    for col in ws.columns:\n"
                "        max_len = 0\n"
                "        col_letter = get_column_letter(col[0].column)\n"
                "        for cell in col:\n"
                "            if cell.value:\n"
                "                max_len = max(max_len, len(str(cell.value)))\n"
                "        ws.column_dimensions[col_letter].width = min(max_len + 2, 60)\n"
            ),
        },
    },
    {
        "name": "pdf",
        "description": (
            "PDF 文档处理技能。支持读取、创建、合并、拆分、提取文字/表格、"
            "添加水印、旋转页面、加密解密、OCR 识别扫描件。"
        ),
        "category": "document",
        "version": "1.0.0",
        "config": {
            "entry": "skill.py",
            "code": (
                "# pdf skill — PDF 文档处理\n"
                "# 依赖: PyPDF2, pdfplumber, reportlab\n"
                "# 安装: pip install PyPDF2 pdfplumber reportlab\n\n"
                "from PyPDF2 import PdfReader, PdfWriter, PdfMerger\n"
                "from reportlab.lib.pagesizes import A4\n"
                "from reportlab.pdfgen import canvas\n"
                "import pdfplumber\n\n"
                "def read_pdf(path: str) -> str:\n"
                "    reader = PdfReader(path)\n"
                "    text = []\n"
                "    for i, page in enumerate(reader.pages):\n"
                "        text.append(f'--- Page {i+1} ---')\n"
                "        text.append(page.extract_text() or '')\n"
                "    return '\\n'.join(text)\n\n"
                "def extract_tables(path: str) -> list:\n"
                "    tables = []\n"
                "    with pdfplumber.open(path) as pdf:\n"
                "        for page in pdf.pages:\n"
                "            t = page.extract_tables()\n"
                "            if t:\n"
                "                tables.extend(t)\n"
                "    return tables\n\n"
                "def merge_pdfs(input_paths: list, output_path: str):\n"
                "    merger = PdfMerger()\n"
                "    for p in input_paths:\n"
                "        merger.append(p)\n"
                "    merger.write(output_path)\n"
                "    merger.close()\n\n"
                "def split_pdf(input_path: str, output_dir: str, pages_per_file: int = 1):\n"
                "    reader = PdfReader(input_path)\n"
                "    total = len(reader.pages)\n"
                "    files = []\n"
                "    for i in range(0, total, pages_per_file):\n"
                "        writer = PdfWriter()\n"
                "        end = min(i + pages_per_file, total)\n"
                "        for j in range(i, end):\n"
                "            writer.add_page(reader.pages[j])\n"
                "        out = f'{output_dir}/split_{i//pages_per_file+1}.pdf'\n"
                "        with open(out, 'wb') as f:\n"
                "            writer.write(f)\n"
                "        files.append(out)\n"
                "    return files\n\n"
                "def add_watermark(input_path: str, output_path: str, text: str):\n"
                "    reader = PdfReader(input_path)\n"
                "    writer = PdfWriter()\n"
                "    for page in reader.pages:\n"
                "        page.merge_page(_create_watermark_page(text))\n"
                "        writer.add_page(page)\n"
                "    with open(output_path, 'wb') as f:\n"
                "        writer.write(f)\n\n"
                "def create_pdf(output_path: str, content: str):\n"
                "    c = canvas.Canvas(output_path, pagesize=A4)\n"
                "    y = 800\n"
                "    for line in content.split('\\n'):\n"
                "        c.drawString(50, y, line[:100])\n"
                "        y -= 15\n"
                "        if y < 50:\n"
                "            c.showPage()\n"
                "            y = 800\n"
                "    c.save()\n"
            ),
        },
    },
    {
        "name": "ppt",
        "description": (
            "PPT 演示文稿处理技能。支持创建、编辑 .pptx 文件，包括幻灯片布局、"
            "文本框、图片、图表、表格、动画、母版、演讲者备注。"
        ),
        "category": "document",
        "version": "1.0.0",
        "config": {
            "entry": "skill.py",
            "code": (
                "# ppt skill — PPT 演示文稿处理\n"
                "# 依赖: python-pptx\n"
                "# 安装: pip install python-pptx\n\n"
                "from pptx import Presentation\n"
                "from pptx.util import Inches, Pt, Emu\n"
                "from pptx.dml.color import RGBColor\n"
                "from pptx.enum.text import PP_ALIGN, MSO_ANCHOR\n\n"
                "def create_presentation():\n"
                "    return Presentation()\n\n"
                "def read_ppt(path: str):\n"
                "    prs = Presentation(path)\n"
                "    content = []\n"
                "    for i, slide in enumerate(prs.slides):\n"
                "        content.append(f'--- Slide {i+1} ---')\n"
                "        for shape in slide.shapes:\n"
                "            if shape.has_text_frame:\n"
                "                content.append(shape.text_frame.text)\n"
                "    return '\\n'.join(content)\n\n"
                "def add_slide(prs, layout_idx: int = 0):\n"
                "    layout = prs.slide_layouts[layout_idx]\n"
                "    return prs.slides.add_slide(layout)\n\n"
                "def add_text_box(slide, text: str, left: float, top: float, width: float, height: float):\n"
                "    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))\n"
                "    tf = txBox.text_frame\n"
                "    tf.text = text\n"
                "    return tf\n\n"
                "def add_title(slide, title: str, subtitle: str = ''):\n"
                "    slide.shapes.title.text = title\n"
                "    if subtitle and len(slide.placeholders) > 1:\n"
                "        slide.placeholders[1].text = subtitle\n\n"
                "def add_image(slide, image_path: str, left: float, top: float, width: float = 5.0):\n"
                "    slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width))\n\n"
                "def add_table(slide, rows: int, cols: int, data: list, left: float, top: float, width: float, height: float):\n"
                "    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))\n"
                "    table = table_shape.table\n"
                "    for i, row_data in enumerate(data):\n"
                "        for j, cell_text in enumerate(row_data):\n"
                "            table.cell(i, j).text = str(cell_text)\n"
                "    return table\n"
            ),
        },
    },
    {
        "name": "web-search",
        "description": (
            "网页搜索技能。支持搜索互联网获取最新信息，包括搜索引擎查询、"
            "网页内容抓取、新闻/百科/技术文档检索。需要联网环境。"
        ),
        "category": "search",
        "version": "1.0.0",
        "config": {
            "entry": "skill.py",
            "code": (
                "# web-search skill — 网页搜索与内容获取\n"
                "# 依赖: requests, beautifulsoup4\n"
                "# 安装: pip install requests beautifulsoup4\n\n"
                "import requests\n"
                "from bs4 import BeautifulSoup\n"
                "from urllib.parse import quote_plus\n\n"
                "USER_AGENT = 'Mozilla/5.0 (compatible; AI-Hubs/4.0; +http://ai-hubs.local)'\n\n"
                "def search_web(query: str, num_results: int = 10) -> list:\n"
                "    \"\"\"使用 DuckDuckGo HTML 搜索（免 API Key）\"\"\"\n"
                "    url = f'https://html.duckduckgo.com/html/?q={quote_plus(query)}'\n"
                "    headers = {'User-Agent': USER_AGENT}\n"
                "    resp = requests.get(url, headers=headers, timeout=15)\n"
                "    soup = BeautifulSoup(resp.text, 'html.parser')\n"
                "    results = []\n"
                "    for item in soup.select('.result')[:num_results]:\n"
                "        title_el = item.select_one('.result__title')\n"
                "        snippet_el = item.select_one('.result__snippet')\n"
                "        link_el = item.select_one('.result__url')\n"
                "        if title_el:\n"
                "            results.append({\n"
                "                'title': title_el.get_text(strip=True),\n"
                "                'snippet': snippet_el.get_text(strip=True) if snippet_el else '',\n"
                "                'url': link_el.get_text(strip=True) if link_el else '',\n"
                "            })\n"
                "    return results\n\n"
                "def fetch_page(url: str, timeout: int = 15) -> str:\n"
                "    \"\"\"抓取网页正文文本\"\"\"\n"
                "    headers = {'User-Agent': USER_AGENT}\n"
                "    resp = requests.get(url, headers=headers, timeout=timeout)\n"
                "    resp.raise_for_status()\n"
                "    soup = BeautifulSoup(resp.text, 'html.parser')\n"
                "    for tag in soup(['script', 'style', 'nav', 'footer', 'header']):\n"
                "        tag.decompose()\n"
                "    text = soup.get_text(separator='\\n', strip=True)\n"
                "    lines = [l.strip() for l in text.split('\\n') if l.strip()]\n"
                "    return '\\n'.join(lines[:200])\n\n"
                "def search_news(query: str, num_results: int = 10) -> list:\n"
                "    \"\"\"搜索新闻\"\"\"\n"
                "    return search_web(f'{query} news', num_results)\n\n"
                "def search_docs(query: str, site: str = '', num_results: int = 10) -> list:\n"
                "    \"\"\"搜索技术文档（可限定站点）\"\"\"\n"
                "    q = f'{query} site:{site}' if site else query\n"
                "    return search_web(q, num_results)\n"
            ),
        },
    },
]


# ═══════════════════════════════════════════════════════════
# 初始化函数
# ═══════════════════════════════════════════════════════════

async def ensure_builtin_skills(session: AsyncSession) -> None:
    """确保 5 个内置技能已存在于数据库（幂等：同名不重复插入）"""
    for skill_def in _BUILTIN_SKILLS:
        existing = (
            await session.execute(
                select(Skill).where(
                    Skill.name == skill_def["name"],
                    Skill.source == "builtin",
                )
            )
        ).scalar_one_or_none()

        if existing is not None:
            continue  # 已存在，跳过

        skill = Skill(
            name=skill_def["name"],
            description=skill_def["description"],
            category=skill_def["category"],
            source="builtin",
            version=skill_def["version"],
            config=skill_def["config"],
            is_installed=True,
            installed_at=datetime.utcnow(),
        )
        session.add(skill)
        logger.info(f"已初始化内置技能: {skill_def['name']}")

    # 确保所有内置技能都已安装
    names = [s["name"] for s in _BUILTIN_SKILLS]
    result = await session.execute(
        select(Skill).where(Skill.name.in_(names), Skill.source == "builtin")
    )
    for sk in result.scalars().all():
        if not sk.is_installed:
            sk.is_installed = True
            sk.installed_at = datetime.utcnow()
            logger.info(f"标记内置技能为已安装: {sk.name}")
